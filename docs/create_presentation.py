#!/usr/bin/env python3
"""
Script to generate PowerPoint presentation for Project Cost & Quality Management System
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor as RgbColor

# Create presentation with 16:9 aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
PRIMARY_COLOR = RgbColor(41, 98, 255)  # Blue
SECONDARY_COLOR = RgbColor(99, 102, 241)  # Indigo
SUCCESS_COLOR = RgbColor(34, 197, 94)  # Green
WARNING_COLOR = RgbColor(234, 179, 8)  # Yellow
DANGER_COLOR = RgbColor(239, 68, 68)  # Red
DARK_COLOR = RgbColor(31, 41, 55)  # Dark gray
LIGHT_COLOR = RgbColor(249, 250, 251)  # Light gray


def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = PRIMARY_COLOR
    background.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = RgbColor(219, 234, 254)
    p.alignment = PP_ALIGN.CENTER

    return slide


def add_section_slide(prs, title, subtitle=""):
    """Add a section divider slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Accent bar on left
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.3), prs.slide_height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY_COLOR
    bar.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(12), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = DARK_COLOR

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(12), Inches(0.8))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(20)
        p.font.color.rgb = RgbColor(107, 114, 128)

    return slide


def add_content_slide(prs, title, content_items, icon=""):
    """Add a content slide with bullet points"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_COLOR
    header.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{icon} {title}" if icon else title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, item in enumerate(content_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_COLOR
        p.space_after = Pt(12)

    return slide


def add_table_slide(prs, title, headers, rows, icon=""):
    """Add a slide with a table"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_COLOR
    header.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{icon} {title}" if icon else title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)

    # Table
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table_width = Inches(12.333)
    table_height = Inches(0.5) * num_rows

    table = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(0.5), Inches(1.6),
        table_width, min(table_height, Inches(5))
    ).table

    # Set column widths
    col_width = table_width / num_cols
    for i in range(num_cols):
        table.columns[i].width = int(col_width)

    # Header row
    for i, header_text in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header_text
        cell.fill.solid()
        cell.fill.fore_color.rgb = SECONDARY_COLOR
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = RgbColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

    # Data rows
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_text)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.font.color.rgb = DARK_COLOR
            # Alternate row colors
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_COLOR

    return slide


def add_metrics_slide(prs, title, metrics):
    """Add a slide with metric cards"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_COLOR
    header.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)

    # Metric cards (2 rows x 3 cols)
    card_width = Inches(3.8)
    card_height = Inches(2.2)
    start_x = Inches(0.5)
    start_y = Inches(1.6)
    gap_x = Inches(0.3)
    gap_y = Inches(0.3)

    for i, metric in enumerate(metrics[:6]):
        row = i // 3
        col = i % 3
        x = start_x + col * (card_width + gap_x)
        y = start_y + row * (card_height + gap_y)

        # Card background
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_width, card_height)
        card.fill.solid()
        card.fill.fore_color.rgb = RgbColor(243, 244, 246)
        card.line.color.rgb = RgbColor(209, 213, 219)

        # Metric name
        name_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.2), card_width - Inches(0.4), Inches(0.4))
        tf = name_box.text_frame
        p = tf.paragraphs[0]
        p.text = metric['name']
        p.font.size = Pt(12)
        p.font.color.rgb = RgbColor(107, 114, 128)

        # Metric value
        value_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.6), card_width - Inches(0.4), Inches(0.8))
        tf = value_box.text_frame
        p = tf.paragraphs[0]
        p.text = metric['value']
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = metric.get('color', DARK_COLOR)

        # Metric description
        desc_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(1.5), card_width - Inches(0.4), Inches(0.5))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = metric['desc']
        p.font.size = Pt(11)
        p.font.color.rgb = RgbColor(107, 114, 128)

    return slide


def add_highlight_slide(prs, title, main_text, sub_text=""):
    """Add a highlight/quote slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Background gradient simulation with shape
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RgbColor(238, 242, 255)  # Light blue
    bg.line.fill.background()

    # Accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6), Inches(2.5), Inches(0.15), Inches(2.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY_COLOR
    bar.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR

    # Main text
    main_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(11.5), Inches(2))
    tf = main_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f'"{main_text}"'
    p.font.size = Pt(26)
    p.font.italic = True
    p.font.color.rgb = DARK_COLOR
    p.alignment = PP_ALIGN.CENTER

    if sub_text:
        sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(5.5), Inches(11.5), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = sub_text
        p.font.size = Pt(16)
        p.font.color.rgb = RgbColor(107, 114, 128)
        p.alignment = PP_ALIGN.CENTER

    return slide


# ============================================
# CREATE PRESENTATION SLIDES
# ============================================

# Slide 1: Title
add_title_slide(
    prs,
    "HỆ THỐNG QUẢN LÝ CHI PHÍ\n& CHẤT LƯỢNG DỰ ÁN",
    "Báo cáo Tổng quan Chức năng | Tháng 01/2026"
)

# Slide 2: Agenda
add_content_slide(prs, "NỘI DUNG TRÌNH BÀY", [
    "Mục tiêu và vấn đề cần giải quyết",
    "Các chức năng đã hoàn thành (15+ modules)",
    "Chỉ số EVM - Dự báo chi phí dự án",
    "Giá trị kinh doanh & ROI",
    "Hướng phát triển tiếp theo"
], "📋")

# Slide 3: Problems
add_section_slide(prs, "VẤN ĐỀ HIỆN TẠI", "Những khó khăn trong quản lý dự án")

# Slide 4: Problem details
add_table_slide(prs, "VẤN ĐỀ & GIẢI PHÁP",
    ["Vấn đề hiện tại", "Giải pháp của hệ thống"],
    [
        ["Không biết dự án đang tốn bao nhiêu", "Dashboard hiển thị chi phí real-time"],
        ["Không dự báo được tổng chi phí", "Hệ thống tự tính EAC (Estimate at Completion)"],
        ["Phát hiện vấn đề khi đã quá muộn", "Cảnh báo tự động: Tốt / Cảnh báo / Rủi ro"],
        ["Quyết định dựa trên cảm tính", "Số liệu và biểu đồ trực quan"],
        ["Mất nhiều thời gian làm báo cáo", "Báo cáo tự động trong vài giây"],
    ],
    "⚠️"
)

# Slide 5: Features section
add_section_slide(prs, "CÁC CHỨC NĂNG ĐÃ HOÀN THÀNH", "15+ modules hoạt động")

# Slide 6: Feature overview
add_table_slide(prs, "TỔNG QUAN CHỨC NĂNG",
    ["Module", "Mô tả", "Lợi ích"],
    [
        ["Dashboard", "Tổng quan tất cả dự án", "Nắm bắt tình hình trong 30 giây"],
        ["Quản lý Phase", "Chia dự án thành giai đoạn", "Phát hiện sớm vấn đề"],
        ["Screen/Function", "Theo dõi từng chức năng", "Không bỏ sót công việc"],
        ["Nhân sự", "Quản lý team & phân công", "Phân bổ nguồn lực hợp lý"],
        ["Effort Tracking", "Ghi nhận công sức", "Số liệu minh bạch"],
        ["Testing/QA", "Theo dõi chất lượng", "Đảm bảo quality"],
    ],
    "🔧"
)

# Slide 7: More features
add_table_slide(prs, "TỔNG QUAN CHỨC NĂNG (tiếp)",
    ["Module", "Mô tả", "Lợi ích"],
    [
        ["Báo cáo tự động", "Tuần / Phase / Project", "Tiết kiệm hàng giờ"],
        ["Chỉ số EVM", "SPI, CPI, EAC, VAC", "Dự báo chi phí chính xác"],
        ["Biểu đồ", "6+ loại charts", "Trực quan dễ hiểu"],
        ["Năng suất Team", "Phân tích theo người/vai trò", "So sánh hiệu suất"],
        ["Cấu hình", "Ngày lễ, giờ làm việc", "Linh hoạt theo công ty"],
    ],
    "🔧"
)

# Slide 8: EVM Section
add_section_slide(prs, "CHỈ SỐ EVM", "Earned Value Management - Chuẩn quốc tế PMI")

# Slide 9: EVM metrics
add_metrics_slide(prs, "CÁC CHỈ SỐ EVM CƠ BẢN", [
    {'name': 'BAC (Budget at Completion)', 'value': '100 MM', 'desc': 'Tổng ngân sách dự kiến', 'color': PRIMARY_COLOR},
    {'name': 'PV (Planned Value)', 'value': '60 MM', 'desc': 'Giá trị kế hoạch đến nay', 'color': SECONDARY_COLOR},
    {'name': 'EV (Earned Value)', 'value': '50 MM', 'desc': 'Giá trị thực tế hoàn thành', 'color': SUCCESS_COLOR},
    {'name': 'AC (Actual Cost)', 'value': '55 MM', 'desc': 'Chi phí thực tế đã bỏ ra', 'color': WARNING_COLOR},
    {'name': 'SPI = EV/PV', 'value': '0.83', 'desc': 'Chậm 17% so với kế hoạch', 'color': DANGER_COLOR},
    {'name': 'CPI = EV/AC', 'value': '0.91', 'desc': 'Vượt 9% ngân sách', 'color': WARNING_COLOR},
])

# Slide 10: Forecasting
add_metrics_slide(prs, "DỰ BÁO CHI PHÍ - \"Cuối cùng tốn bao nhiêu?\"", [
    {'name': 'EAC (Estimate at Completion)', 'value': '110 MM', 'desc': 'Dự kiến tổng chi phí', 'color': DANGER_COLOR},
    {'name': 'VAC (Variance at Completion)', 'value': '-10 MM', 'desc': 'Vượt 10 MM so với budget', 'color': DANGER_COLOR},
    {'name': 'TCPI', 'value': '1.11', 'desc': 'Cần tăng 11% hiệu suất', 'color': WARNING_COLOR},
    {'name': 'Tiến độ', 'value': '50%', 'desc': 'Đã hoàn thành', 'color': PRIMARY_COLOR},
    {'name': 'Budget còn lại', 'value': '45 MM', 'desc': 'Để hoàn thành 50% còn lại', 'color': SECONDARY_COLOR},
    {'name': 'Trạng thái', 'value': 'Cảnh báo', 'desc': 'Cần review kế hoạch', 'color': WARNING_COLOR},
])

# Slide 11: Status explanation
add_table_slide(prs, "Ý NGHĨA TRẠNG THÁI DỰ ÁN",
    ["Trạng thái", "Điều kiện", "Hành động"],
    [
        ["🟢 GOOD", "CPI ≥ 1.0, Pass Rate ≥ 95%", "Tiếp tục theo dõi"],
        ["🟡 WARNING", "CPI 0.83-1.0 hoặc Pass Rate 80-95%", "Review kế hoạch"],
        ["🔴 AT RISK", "CPI < 0.83 hoặc Pass Rate < 80%", "Can thiệp ngay"],
    ],
    "🚦"
)

# Slide 12: Business Value
add_section_slide(prs, "GIÁ TRỊ KINH DOANH", "ROI và lợi ích đạt được")

# Slide 13: Time savings
add_table_slide(prs, "TIẾT KIỆM THỜI GIAN",
    ["Công việc", "Trước đây", "Với hệ thống", "Tiết kiệm"],
    [
        ["Tổng hợp báo cáo tuần", "2-4 giờ", "5 phút", "~95%"],
        ["Kiểm tra tình trạng dự án", "Họp 1 giờ", "30 giây", "~99%"],
        ["Tính dự báo ngân sách", "Nửa ngày", "Tự động", "100%"],
        ["Setup dự án mới", "1-2 ngày", "30 phút", "~90%"],
    ],
    "⏱️"
)

# Slide 14: ROI
add_content_slide(prs, "ROI ƯỚC TÍNH", [
    "5 PM sử dụng hệ thống × 4 giờ tiết kiệm/tuần = 20 giờ/tuần",
    "20 giờ × 4 tuần × 12 tháng = 960 giờ/năm",
    "Quy đổi: ~5.5 man-month/năm tiết kiệm được",
    "",
    "GIÁ TRỊ KHÁC (khó đo lường):",
    "• Phát hiện sớm dự án có vấn đề → Tiết kiệm chi phí sửa chữa",
    "• Báo cáo chuyên nghiệp → Tăng độ tin cậy với khách hàng",
    "• Dữ liệu lịch sử → Cải thiện ước lượng dự án tương lai"
], "💰")

# Slide 15: Risk reduction
add_table_slide(prs, "GIẢM THIỂU RỦI RO",
    ["Rủi ro", "Cách hệ thống giúp giảm thiểu"],
    [
        ["Vượt ngân sách", "Cảnh báo sớm khi CPI < 1.0, dự báo EAC"],
        ["Trễ deadline", "Theo dõi SPI, delay rate, cảnh báo khi chậm"],
        ["Chất lượng kém", "Theo dõi pass rate, defect density"],
        ["Thiếu minh bạch", "Dữ liệu real-time, báo cáo tự động"],
        ["Quyết định sai", "Ra quyết định dựa trên dữ liệu"],
    ],
    "🛡️"
)

# Slide 16: Future roadmap
add_section_slide(prs, "HƯỚNG PHÁT TRIỂN", "Giai đoạn tiếp theo")

# Slide 17: Roadmap details
add_table_slide(prs, "ĐỀ XUẤT PHÁT TRIỂN GIAI ĐOẠN 2",
    ["Tính năng", "Giá trị kỳ vọng", "Ưu tiên"],
    [
        ["Tích hợp Jira/Azure DevOps", "Đồng bộ dữ liệu tự động", "Cao"],
        ["Dashboard cho khách hàng", "Khách tự theo dõi, tăng tin cậy", "Cao"],
        ["Xuất PDF/Excel", "Gửi báo cáo cho stakeholders", "Trung bình"],
        ["Mobile App", "Xem báo cáo mọi lúc mọi nơi", "Trung bình"],
        ["So sánh nhiều dự án", "Benchmark hiệu suất", "Thấp"],
    ],
    "🚀"
)

# Slide 18: Summary
add_section_slide(prs, "TÓM TẮT", "")

# Slide 19: Summary content
add_metrics_slide(prs, "NHỮNG GÌ ĐÃ HOÀN THÀNH", [
    {'name': 'Modules chức năng', 'value': '15+', 'desc': 'Đã hoạt động', 'color': SUCCESS_COLOR},
    {'name': 'API Endpoints', 'value': '60+', 'desc': 'Backend services', 'color': PRIMARY_COLOR},
    {'name': 'Màn hình giao diện', 'value': '10+', 'desc': 'Frontend screens', 'color': SECONDARY_COLOR},
    {'name': 'Loại biểu đồ', 'value': '6', 'desc': 'Charts & Graphs', 'color': PRIMARY_COLOR},
    {'name': 'Loại báo cáo', 'value': '3', 'desc': 'Tuần/Phase/Project', 'color': SECONDARY_COLOR},
    {'name': 'Tự động hóa', 'value': '80%', 'desc': 'Công việc báo cáo', 'color': SUCCESS_COLOR},
])

# Slide 20: Key message
add_highlight_slide(
    prs,
    "THÔNG ĐIỆP CHÍNH",
    "Hệ thống không chỉ THEO DÕI dự án - mà còn DỰ BÁO và CẢNH BÁO SỚM để ban lãnh đạo có thể HÀNH ĐỘNG kịp thời, tiết kiệm chi phí và giảm rủi ro.",
    "Áp dụng chuẩn EVM quốc tế | Tự động hóa 80% báo cáo | Trực quan hóa dữ liệu"
)

# Slide 21: Thank you
add_title_slide(
    prs,
    "CẢM ƠN QUÝ VỊ\nĐÃ LẮNG NGHE",
    "Sẵn sàng giải đáp thắc mắc"
)

# Save presentation
output_path = '/home/user/project-cost-quality-management/docs/PROJECT_OVERVIEW_PRESENTATION.pptx'
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
